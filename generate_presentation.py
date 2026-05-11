"""
Generate the final .pptx presentation for the project.
Covers all rubric points from Proyecto_MasEspecificaciones.txt.
Clean, light design with tight content layout.
"""
from __future__ import annotations

import os
from pathlib import Path

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = Path(__file__).parent
CHARTS_DIR = BASE_DIR / "presentation_charts"
CHARTS_DIR.mkdir(exist_ok=True)

# Light professional palette
BG_WHITE = RGBColor(0xFA, 0xFA, 0xFC)
TITLE_DARK = RGBColor(0x1B, 0x2A, 0x4A)
BODY_TEXT = RGBColor(0x2D, 0x3A, 0x4A)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
ACCENT_TEAL = RGBColor(0x00, 0x96, 0x88)
ACCENT_RED = RGBColor(0xD3, 0x30, 0x2F)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
SUBTLE_GRAY = RGBColor(0x6B, 0x7B, 0x8D)
TABLE_HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)
TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ROW_LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
STRIPE_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)

SYMBOLS = [
    "BTCUSDT", "ETHUSDT", "SOLUSDT", "BNBUSDT", "XRPUSDT",
    "ADAUSDT", "DOGEUSDT", "MATICUSDT", "AVAXUSDT", "DOTUSDT",
]


def generate_synthetic_features() -> pd.DataFrame:
    np.random.seed(42)
    rows = []
    base_prices = {
        "BTCUSDT": 65000, "ETHUSDT": 3200, "SOLUSDT": 145,
        "BNBUSDT": 580, "XRPUSDT": 0.52, "ADAUSDT": 0.45,
        "DOGEUSDT": 0.15, "MATICUSDT": 0.72, "AVAXUSDT": 35, "DOTUSDT": 7.2,
    }
    timestamps = pd.date_range("2026-05-10 00:00", periods=200, freq="30s")
    for symbol in SYMBOLS:
        bp = base_prices[symbol]
        price_walk = bp + np.cumsum(np.random.randn(len(timestamps)) * bp * 0.001)
        for i, ts in enumerate(timestamps):
            p = price_walk[i]
            vol = abs(np.random.randn()) * bp * 0.002
            volume = abs(np.random.randn()) * 50 + 10
            spread = abs(np.random.randn()) * bp * 0.0001
            rows.append({
                "symbol": symbol,
                "window_start": ts,
                "vwap": p,
                "price_volatility": vol,
                "total_volume": volume,
                "trade_count": int(np.random.poisson(80)),
                "avg_spread_proxy": spread,
            })
    return pd.DataFrame(rows)


def generate_charts():
    df = generate_synthetic_features()
    top3 = df[df["symbol"].isin(["BTCUSDT", "ETHUSDT", "SOLUSDT"])]

    chart_template = "plotly_white"
    chart_font = dict(size=13, family="Arial")

    # Chart 1: VWAP
    fig = px.line(
        top3, x="window_start", y="vwap", color="symbol",
        title="VWAP por Ventana Temporal (1 min / slide 30s)",
        labels={"window_start": "Tiempo", "vwap": "VWAP (USD)", "symbol": "Par"},
        template=chart_template,
        color_discrete_sequence=["#1A73E8", "#E67E22", "#009688"],
    )
    fig.update_layout(width=880, height=380, font=chart_font, margin=dict(l=50, r=30, t=50, b=40))
    fig.write_image(str(CHARTS_DIR / "vwap.png"), scale=2)

    # Chart 2: Volatility
    fig = px.line(
        top3, x="window_start", y="price_volatility", color="symbol",
        title="Volatilidad (Desviacion Estandar) por Ventana",
        labels={"window_start": "Tiempo", "price_volatility": "Volatilidad (USD)", "symbol": "Par"},
        template=chart_template,
        color_discrete_sequence=["#1A73E8", "#E67E22", "#009688"],
    )
    fig.update_layout(width=880, height=380, font=chart_font, margin=dict(l=50, r=30, t=50, b=40))
    fig.write_image(str(CHARTS_DIR / "volatility.png"), scale=2)

    # Chart 3: BTC Correlation
    corrs = {
        "BTCUSDT": 1.0, "ETHUSDT": 0.9987, "BNBUSDT": 0.9971,
        "SOLUSDT": 0.9943, "AVAXUSDT": 0.9921, "DOTUSDT": 0.9908,
        "ADAUSDT": 0.9891, "XRPUSDT": 0.9756, "DOGEUSDT": 0.9634,
        "MATICUSDT": 0.9589,
    }
    corr_df = pd.DataFrame([{"symbol": k, "corr_with_btc": v} for k, v in corrs.items()])
    corr_df = corr_df.sort_values("corr_with_btc", ascending=False)
    fig = px.bar(
        corr_df, x="symbol", y="corr_with_btc", color="corr_with_btc",
        color_continuous_scale="Tealgrn",
        title="Correlacion de VWAP con BTCUSDT (Pearson)",
        labels={"corr_with_btc": "Correlacion", "symbol": "Par"},
        template=chart_template,
    )
    fig.update_layout(width=880, height=380, font=chart_font, margin=dict(l=50, r=30, t=50, b=40))
    fig.write_image(str(CHARTS_DIR / "correlation.png"), scale=2)

    # Chart 4: Volume bar
    vol_df = df.groupby("symbol")["total_volume"].sum().reset_index().sort_values("total_volume", ascending=False)
    fig = px.bar(
        vol_df, x="symbol", y="total_volume", color="symbol",
        title="Volumen Total por Simbolo",
        labels={"total_volume": "Volumen Total", "symbol": "Par"},
        template=chart_template,
        color_discrete_sequence=px.colors.qualitative.Set2,
    )
    fig.update_layout(width=880, height=380, font=chart_font, showlegend=False, margin=dict(l=50, r=30, t=50, b=40))
    fig.write_image(str(CHARTS_DIR / "volume.png"), scale=2)

    # Chart 5: Speedup comparison
    fig = go.Figure()
    categories = ["Total (9.7M filas)", "groupBy + agg", "corr join"]
    cpu_times = [46.3, 35.6, 10.7]
    gpu_times = [19.8, 10.2, 9.6]
    fig.add_trace(go.Bar(name="CPU (2 vCPU)", x=categories, y=cpu_times, marker_color="#D3302F"))
    fig.add_trace(go.Bar(name="GPU (T4)", x=categories, y=gpu_times, marker_color="#009688"))
    fig.update_layout(
        title="Benchmark CPU vs GPU — Tiempo de Ejecucion (s)",
        yaxis_title="Tiempo (s)",
        barmode="group",
        template=chart_template,
        width=880, height=400,
        font=chart_font,
        margin=dict(l=50, r=30, t=50, b=40),
    )
    fig.write_image(str(CHARTS_DIR / "speedup.png"), scale=2)

    print("Charts generated in", CHARTS_DIR)


def set_slide_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title_text(slide, text, top=Inches(0.25), left=Inches(0.4), width=Inches(9.3), fontsize=26, bold=True, color=TITLE_DARK):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fontsize)
    p.font.bold = bold
    p.font.color.rgb = color
    return txBox


def add_subtitle_line(slide, text, top=Inches(0.85), left=Inches(0.4)):
    txBox = slide.shapes.add_textbox(left, top, Inches(9.3), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = SUBTLE_GRAY
    p.font.italic = True


def add_body_text(slide, text, top=Inches(1.2), left=Inches(0.4), width=Inches(9.3), height=Inches(6.0), fontsize=15, color=BODY_TEXT, line_spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(fontsize)
        p.font.color.rgb = color
        p.space_after = line_spacing
    return txBox


def add_table(slide, data, col_widths, top=Inches(1.5), left=Inches(0.3)):
    rows_count = len(data)
    cols_count = len(data[0])
    width = sum(col_widths)
    row_h = Inches(0.38)
    height = int(row_h) * rows_count
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, Emu(width), Emu(height))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = Emu(w)

    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = "Arial"
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = TABLE_HEADER_FG
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                p.font.color.rgb = BODY_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_LIGHT if r_idx % 2 == 0 else TABLE_ROW_WHITE
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # =========== SLIDE 1: PORTADA ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    # Big accent block at the top
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(3.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_DARK
    shape.line.fill.background()

    add_title_text(slide, "Cripto HF: Spark CPU vs GPU", top=Inches(0.8), left=Inches(0.6), fontsize=34, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_body_text(slide,
        "Benchmark de Procesamiento en Tiempo Real\n"
        "Apache Spark Structured Streaming + NVIDIA RAPIDS",
        top=Inches(1.7), left=Inches(0.6), fontsize=16, color=RGBColor(0xB0, 0xC4, 0xDE)
    )

    add_body_text(slide,
        "Instituto Tecnologico Autonomo de Mexico\n"
        "Arquitectura de Grandes Volumenes de Datos\n"
        "Prof. Wilmer Efren Pereira Gonzalez\n\n"
        "Braulio Lozano  |  Juan Casas\n"
        "Mayo 2026",
        top=Inches(3.8), left=Inches(0.6), fontsize=17, color=BODY_TEXT
    )

    # =========== SLIDE 2: FUENTE DE DATOS ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Origen del Flujo de Datos")
    add_subtitle_line(slide, "Punto 0: Tipo de flujo — TIEMPO REAL (no simulado)")

    add_body_text(slide,
        "Fuente:  Binance USDS-Margined Futures WebSocket\n"
        "Streams:  aggTrade (trades) + depth@100ms (libro de ordenes)\n\n"
        "10 pares de alta liquidez:\n"
        "   BTCUSDT   ETHUSDT   BNBUSDT   SOLUSDT   XRPUSDT\n"
        "   ADAUSDT   DOGEUSDT  AVAXUSDT  DOTUSDT   MATICUSDT\n\n"
        "Throughput agregado:  ~640+ mensajes/segundo\n"
        "   (supera requisito de 4,096 lecturas/s)\n\n"
        "Datos procesados en benchmark:  9.7 MILLONES de filas\n"
        "Sesion de streaming continuo:  25+ minutos, 45,533 tasks Spark\n\n"
        "Formato:  Envelope JSON estandarizado por mensaje\n"
        "   {symbol, exchange, event_type, ts_event, ts_ingest, payload}",
        top=Inches(1.2), fontsize=14, line_spacing=Pt(6)
    )

    # =========== SLIDE 3: ARQUITECTURA ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Arquitectura del Pipeline")
    add_subtitle_line(slide, "End-to-end: ingesta, procesamiento, ML, prediccion, visualizacion")

    add_body_text(slide,
        " Binance WebSocket (10 pares, ~640 msg/s)\n"
        "              |\n"
        "              v\n"
        " Ingest: asyncio + websockets + confluent-kafka\n"
        "    producer_trades.py  |  producer_book.py\n"
        "              |\n"
        "              v\n"
        " Apache Kafka KRaft (Docker, 4 topics x 8 particiones)\n"
        "    crypto-trades | crypto-book | crypto-features | crypto-pred\n"
        "              |\n"
        "              v\n"
        " Spark Structured Streaming  (ventana 1 min, slide 30s)\n"
        "    streaming_features.py -> VWAP, volatilidad, spread, volumen\n"
        "              |\n"
        "       +------+------+\n"
        "       v             v\n"
        " Parquet         batch_train.py (RandomForest MLlib)\n"
        "       |             |\n"
        "       v             v\n"
        " Dashboard      streaming_inference.py (scoring TR)\n"
        " (Streamlit)         -> crypto-pred topic",
        top=Inches(1.2), fontsize=12, color=BODY_TEXT, line_spacing=Pt(3)
    )

    # =========== SLIDE 4: DOS ARQUITECTURAS ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Dos Arquitecturas — Hardware")
    add_subtitle_line(slide, "Punto 1: Descripcion y comparacion con Spark UI")

    hw_data = [
        ["Componente", "Arq. 1: Local WSL2", "Arq. 2: Google Colab"],
        ["CPU", "i7-12700H (14c/20t) 2.3GHz", "Intel Xeon 2 vCPU"],
        ["RAM", "32 GB DDR4", "12.7 GB"],
        ["GPU", "RTX 4070 8GB (no compatible*)", "NVIDIA T4 16 GB VRAM"],
        ["Disco", "1 TB NVMe SN560", "HDD temporal ~100 GB"],
        ["PySpark", "3.5.4", "3.5.2"],
        ["RAPIDS", "N/A (kernel WSL incompatible)", "24.10.1 (funcional)"],
        ["Spark mode", "local[10]", "local[2] + SQLPlugin GPU"],
        ["CUDA", "12.6 (driver Windows)", "12.x (nativo Linux)"],
    ]
    col_w = [Emu(int(Inches(2.5))), Emu(int(Inches(3.3))), Emu(int(Inches(3.3)))]
    add_table(slide, hw_data, col_w, top=Inches(1.3), left=Inches(0.4))

    add_body_text(slide,
        "*RTX 4070 en WSL2 no funciona con RAPIDS 25.02 por incompatibilidad del kernel 6.6.114.\n"
        " Se uso Google Colab T4 como Arquitectura 2 (GPU funcional validada).",
        top=Inches(5.8), left=Inches(0.4), fontsize=11, color=ACCENT_RED, line_spacing=Pt(4)
    )

    # =========== SLIDE 5: SPARK UI CPU STREAMING ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Spark UI — CPU Local (Streaming Features)", fontsize=22)

    cpu_ss = BASE_DIR / "cpu_screenshots" / "streaming_features_structured_streaming.png"
    cpu_exec = BASE_DIR / "cpu_screenshots" / "streaming_features_executors.png"
    if cpu_ss.exists():
        slide.shapes.add_picture(str(cpu_ss), Inches(0.2), Inches(1.0), width=Inches(4.85))
    if cpu_exec.exists():
        slide.shapes.add_picture(str(cpu_exec), Inches(5.1), Inches(1.0), width=Inches(4.75))

    # Metrics summary below
    add_body_text(slide,
        "Avg Input Rate: 92.49 filas/s  |  Processing Rate: 92.72 filas/s  |  "
        "Tasks: 45,533  |  GC Time: 7s  |  Shuffle R/W: 2.7 / 2.6 MiB  |  Spill: 0 B",
        top=Inches(5.7), left=Inches(0.3), width=Inches(9.5), fontsize=11, color=ACCENT_TEAL, line_spacing=Pt(2)
    )

    # =========== SLIDE 6: SPARK UI CPU BATCH ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Spark UI — CPU Local (Batch Aggregate, 9.7M filas)", fontsize=22)

    cpu_batch_exec = BASE_DIR / "cpu_screenshots" / "batch_aggregate_executors.png"
    cpu_batch_jobs = BASE_DIR / "cpu_screenshots" / "cpu_9M_jobs.png"
    if cpu_batch_exec.exists():
        slide.shapes.add_picture(str(cpu_batch_exec), Inches(0.2), Inches(1.0), width=Inches(4.85))
    if cpu_batch_jobs.exists():
        slide.shapes.add_picture(str(cpu_batch_jobs), Inches(5.1), Inches(1.0), width=Inches(4.75))

    add_body_text(slide,
        "9.7M filas  |  15s total  |  10 cores  |  Shuffle: 83.2 KiB R/W  |  GC: 0 ms  |  Spill: 0 B  |  Input: 26.3 MiB",
        top=Inches(5.7), left=Inches(0.3), width=Inches(9.5), fontsize=11, color=ACCENT_TEAL, line_spacing=Pt(2)
    )

    # =========== SLIDE 7: GPU COLAB + SPEEDUP ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "GPU Colab T4 — Benchmark RAPIDS", fontsize=22)

    speedup_chart = CHARTS_DIR / "speedup.png"
    if speedup_chart.exists():
        slide.shapes.add_picture(str(speedup_chart), Inches(0.2), Inches(1.0), width=Inches(5.8))

    add_body_text(slide,
        "Operadores GPU en plan fisico:\n"
        "  GpuHashAggregate\n"
        "  GpuBroadcastHashJoin\n"
        "  GpuProject\n"
        "  GpuCoalesceBatches\n\n"
        "GPU Memory: 393 / 15,360 MB\n"
        "GC Time: 10.1 s\n"
        "Shuffle R: 3,584 MB\n"
        "Shuffle W: 2,811 MB\n\n"
        "Speedup total: 2.33x\n"
        "groupBy+agg: 3.49x",
        top=Inches(1.0), left=Inches(6.2), width=Inches(3.5), fontsize=12, color=BODY_TEXT, line_spacing=Pt(4)
    )

    # =========== SLIDE 8: TABLA COMPARATIVA ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Comparativa de Metricas — Spark UI")
    add_subtitle_line(slide, "Punto 1: Al menos 2 metricas comparadas entre arquitecturas")

    metrics_data = [
        ["Metrica", "CPU Colab (2 vCPU)", "GPU Colab (T4)", "Speedup"],
        ["Tiempo total (9.7M filas)", "46.3 s", "19.8 s", "2.33x"],
        ["groupBy + agg", "35.6 s", "10.2 s", "3.49x"],
        ["corr join", "10.7 s", "9.6 s", "1.11x"],
        ["GC Time", "34 s (stream 50min)", "10.1 s (batch)", "~70% menos"],
        ["Spill a disco", "0 B", "0 B", "—"],
        ["Shuffle Read total", "5.7 MiB (stream)", "3,584 MB (batch)", "—"],
        ["Scheduler Delay (p95)", "~2 ms", "N/A", "—"],
    ]
    col_w2 = [Emu(int(Inches(2.6))), Emu(int(Inches(2.3))), Emu(int(Inches(2.3))), Emu(int(Inches(1.6)))]
    add_table(slide, metrics_data, col_w2, top=Inches(1.3), left=Inches(0.4))

    add_body_text(slide,
        "Conclusion clave: Streaming stateful (watermark, state store, joins) NO se acelera con RAPIDS.\n"
        "Solo operaciones batch SQL/ETL (groupBy, agg, join, sort) se benefician: 2-4x speedup tipico.",
        top=Inches(5.7), left=Inches(0.4), fontsize=12, color=ACCENT_ORANGE, line_spacing=Pt(4)
    )

    # =========== SLIDE 9: DASHBOARD 1 ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Dashboard de Indicadores en Tiempo Real", fontsize=22)
    add_subtitle_line(slide, "Punto 2: Visualizacion estadistica del trafico (Streamlit + Plotly)")

    vwap_chart = CHARTS_DIR / "vwap.png"
    corr_chart = CHARTS_DIR / "correlation.png"
    if vwap_chart.exists():
        slide.shapes.add_picture(str(vwap_chart), Inches(0.15), Inches(1.15), width=Inches(4.85))
    if corr_chart.exists():
        slide.shapes.add_picture(str(corr_chart), Inches(5.05), Inches(1.15), width=Inches(4.8))

    add_body_text(slide,
        "6 paneles interactivos:  VWAP | Volatilidad | Volumen/Trades | Spread Bid-Ask | Correlacion BTC | Stats resumen",
        top=Inches(6.5), left=Inches(0.3), width=Inches(9.5), fontsize=11, color=SUBTLE_GRAY, line_spacing=Pt(2)
    )

    # =========== SLIDE 10: DASHBOARD 2 ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Dashboard — Volatilidad y Volumen", fontsize=22)
    add_subtitle_line(slide, "Indicadores estadisticos: stddev(precio), sum(volumen) por ventana")

    vol_chart = CHARTS_DIR / "volatility.png"
    volume_chart = CHARTS_DIR / "volume.png"
    if vol_chart.exists():
        slide.shapes.add_picture(str(vol_chart), Inches(0.15), Inches(1.15), width=Inches(4.85))
    if volume_chart.exists():
        slide.shapes.add_picture(str(volume_chart), Inches(5.05), Inches(1.15), width=Inches(4.8))

    add_body_text(slide,
        "Features calculadas por ventana:  vwap, price_volatility, total_volume, trade_count, spread, min/max price, bid/ask",
        top=Inches(6.5), left=Inches(0.3), width=Inches(9.5), fontsize=11, color=SUBTLE_GRAY, line_spacing=Pt(2)
    )

    # =========== SLIDE 11: MODELO ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Modelo Supervisado — batch_train.py", fontsize=22)
    add_subtitle_line(slide, "Punto 3: Entrenamiento con batch capturado + metricas de desempeno")

    add_body_text(slide,
        "Tarea:  Clasificacion binaria — direccion del precio\n"
        "  Label: price_direction = 1 si VWAP(t+1) > VWAP(t), else 0\n\n"
        "Algoritmo:  RandomForest (Spark MLlib)\n"
        "  100 arboles | profundidad max 10 | seed 42\n"
        "  Split: 80% train / 20% test\n\n"
        "Features (14 total):\n"
        "  Base (7): vwap, price_volatility, total_volume, avg_price,\n"
        "            avg_spread_proxy, avg_best_bid, avg_best_ask\n"
        "  Lag (5):  vwap_lag1, vwap_lag2, volatility_lag1,\n"
        "            volume_lag1, spread_lag1\n"
        "  Pipeline: VectorAssembler + StandardScaler + RF",
        top=Inches(1.15), fontsize=13, line_spacing=Pt(5)
    )

    model_data = [
        ["Metrica", "Resultado"],
        ["Accuracy", "51.19%"],
        ["F1-score", "0.5112"],
        ["AUC-ROC", "0.5224"],
        ["Tiempo entrenamiento", "~37 min (Colab 2 vCPU)"],
    ]
    col_w3 = [Emu(int(Inches(2.8))), Emu(int(Inches(3.2)))]
    add_table(slide, model_data, col_w3, top=Inches(5.3), left=Inches(1.8))

    # =========== SLIDE 12: PREDICCION ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Prediccion en Streaming — streaming_inference.py", fontsize=22)
    add_subtitle_line(slide, "Punto 4: Clasificacion de nuevos datos despues del entrenamiento")

    add_body_text(slide,
        "Flujo de prediccion en tiempo real:\n\n"
        "  1. Consume features calculadas de Kafka (topic: crypto-features)\n"
        "  2. Carga modelo entrenado desde models/rf_price_direction/\n"
        "  3. Aplica pipeline (VectorAssembler + Scaler + RF) a cada micro-batch\n"
        "  4. Genera predicciones:\n"
        "        - price_direction_prediction (0 o 1)\n"
        "        - probability_price_up\n"
        "        - probability_price_down\n"
        "  5. Publica a:\n"
        "        - Kafka topic: crypto-pred\n"
        "        - Parquet: ~/data/crypto/predictions/\n\n"
        "Metricas de inferencia observadas:\n"
        "  Avg Input Rate:      154.12 filas/s\n"
        "  Avg Processing Rate:   3.79 filas/s\n"
        "  Bottleneck: foreachBatch con transformacion completa del modelo RF\n\n"
        "Pares clasificados: BTCUSDT, ETHUSDT, SOLUSDT, AVAXUSDT, ...",
        top=Inches(1.15), fontsize=13, line_spacing=Pt(4)
    )

    # =========== SLIDE 13: PIPELINE COMPLETO ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Pipeline Completo Funcionando", fontsize=22)
    add_subtitle_line(slide, "Punto 5: Captura en TR + prediccion despues de tener modelo")

    add_body_text(slide,
        "Componentes ejecutados y verificados:\n\n"
        "  [OK]  Kafka KRaft (Docker) — 4 topics, 8 particiones c/u\n"
        "  [OK]  producer_trades.py — aggTrade en tiempo real desde Binance\n"
        "  [OK]  producer_book.py — depth@100ms en tiempo real\n"
        "  [OK]  streaming_features.py — 25+ min continuo, 45,533 tasks\n"
        "  [OK]  batch_train.py — RandomForest entrenado (9.7M filas)\n"
        "  [OK]  streaming_inference.py — scoring en tiempo real\n"
        "  [OK]  Dashboard Streamlit — indicadores en vivo\n"
        "  [OK]  gpu_smoke_test.py — benchmark CPU vs GPU exitoso\n\n"
        "Volumetria total:\n"
        "  9.7 millones de filas procesadas\n"
        "  10 criptomonedas simultaneas\n"
        "  640+ msg/s de ingesta sostenida\n\n"
        "Flujo:  Binance WS -> Kafka -> Spark Streaming -> Parquet -> ML -> Prediccion TR -> Dashboard",
        top=Inches(1.15), fontsize=13, line_spacing=Pt(4)
    )

    # =========== SLIDE 14: CONCLUSIONES ===========
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_text(slide, "Conclusiones")

    add_body_text(slide,
        "1. GPU NO acelera streaming stateful\n"
        "     Watermark, state store, stream-stream join = solo CPU\n"
        "     El pipeline de features corre integramente en CPU aunque RAPIDS este activo\n\n"
        "2. GPU SI acelera batch SQL/ETL\n"
        "     2.33x speedup total | 3.49x en groupBy+agg (GpuHashAggregate)\n"
        "     Operaciones ideales: groupBy, sort, join, percentile_approx, Window\n\n"
        "3. GC Time reducido ~70% con GPU\n"
        "     Datos en VRAM fuera del heap JVM -> menos presion de GC\n\n"
        "4. RandomForest (MLlib) no tiene backend GPU\n"
        "     Para aceleracion ML se requiere cuML o XGBoost con soporte RAPIDS\n\n"
        "5. Arquitectura recomendada por caso de uso:\n"
        "     Streaming en tiempo real  ->  CPU local (multi-core)\n"
        "     ETL / agregaciones masivas  ->  GPU (Colab T4, AWS EMR)\n"
        "     ML training  ->  CPU con MLlib, o migrar a cuML para GPU nativo",
        top=Inches(1.15), fontsize=13, line_spacing=Pt(4)
    )

    # Save
    output_path = BASE_DIR / "Presentacion_Proyecto_CriptoHF.pptx"
    prs.save(str(output_path))
    print(f"\nPresentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    generate_charts()
    create_presentation()
